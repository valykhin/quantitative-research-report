"""Сборка презентации PPTX по результатам исследования."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

APP = Path(__file__).parent
PLOTS = APP / "output" / "plots"
OUT = APP / "output"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x00, 0x61, 0x00)
RED = RGBColor(0x9C, 0x00, 0x06)
GRAY = RGBColor(0x60, 0x60, 0x60)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BLUE

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Анализ рынка организации\nбеговых мероприятий в России"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Исследование рынка, ориентированное на сбыт"
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf2.add_paragraph()
    p3.text = "\nИсточник: ГИРБО (bo.nalog.gov.ru) | Период: 2021–2025\n11 компаний | 49 наблюдений"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
    p3.alignment = PP_ALIGN.CENTER


def add_slide(prs, title, bullets=None, image_path=None, image_width=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE

    if bullets and image_path:
        # Left: bullets, Right: image
        text_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(4.5), Inches(5.8))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for i, bullet in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = bullet
            para.font.size = Pt(13)
            para.font.color.rgb = DARK
            para.space_after = Pt(8)

        w = image_width or Inches(5.0)
        slide.shapes.add_picture(str(image_path), Inches(5.0), Inches(1.1), width=w)

    elif bullets:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9.0), Inches(5.8))
        tf = text_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if bullet.startswith("!"):
                para.text = bullet[1:]
                para.font.size = Pt(14)
                para.font.bold = True
                para.font.color.rgb = BLUE
                para.space_before = Pt(12)
            elif bullet.startswith("+"):
                para.text = bullet[1:]
                para.font.size = Pt(13)
                para.font.color.rgb = GREEN
                para.font.bold = True
            elif bullet.startswith("-"):
                para.text = bullet[1:]
                para.font.size = Pt(13)
                para.font.color.rgb = RED
            else:
                para.text = bullet
                para.font.size = Pt(13)
                para.font.color.rgb = DARK
            para.space_after = Pt(6)

    elif image_path:
        w = image_width or Inches(9.0)
        left = (Inches(10) - w) / 2
        slide.shapes.add_picture(str(image_path), left, Inches(1.1), width=w)

    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = sum(col_widths) if col_widths else Inches(9.0)
    left = (Inches(10) - tbl_width) / 2
    table = slide.shapes.add_table(n_rows, n_cols, left, Inches(1.2), tbl_width, Inches(0.4 * n_rows)).table

    if col_widths:
        for j, w in enumerate(col_widths):
            table.columns[j].width = w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.color.rgb = DARK
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Титул
    add_title_slide(prs)

    # 2. Цель и гипотезы
    add_slide(prs, "Цель и гипотезы исследования", [
        "!Цель",
        "Оценить финансовое состояние и динамику развития компаний-организаторов беговых мероприятий в РФ (2021–2025)",
        "",
        "!Гипотезы",
        "H1: Рынок демонстрирует устойчивый рост выручки",
        "H2: Существует зависимость между масштабом и рентабельностью",
        "H3: Доля себестоимости снижается с ростом выручки (эффект масштаба)",
        "H4: Компании образуют кластеры по финансовому профилю",
        "H5: Рост выручки сопровождается ростом кредиторской задолженности",
        "",
        "!Метод",
        "Вторичные данные из ГИРБО (API ФНС) | 11 компаний | 49 наблюдений",
    ])

    # 3. Выборка
    add_table_slide(prs, "Состав выборки",
        ["Бренд", "Юрлицо", "Специализация"],
        [
            ["Гонка Героев", "АНО «ЛИГА ГЕРОЕВ»", "Забеги с препятствиями"],
            ["Марафон Сервис", "ООО «МАРАФОН СЕРВИС»", "Марафоны"],
            ["TIMERMAN", "ООО «АСМ «НОВЫЙ СПОРТ»", "Триатлон, забеги"],
            ["Бегом по Золотому кольцу", "ООО «АРЕНА ПЛЮС»", "Беговые туры"],
            ["IRONSTAR", "ООО «АРХИТЕКТУРА СПОРТА»", "Триатлон, забеги"],
            ["ГРУТ / Вайлд Трейл", "ООО «ГРУТ»", "Трейлраннинг"],
            ["parkrun", "АНО «ПАРКРАН»", "Парковые забеги"],
            ["PushkinRun", "АССОЦИАЦИЯ «ПЕТЕРБУРГСКИЙ СПОРТ»", "Городские забеги"],
            ["Беги, Герой!", "АНО «РЕЙТИНГ СПОРТ»", "Городские забеги"],
            ["S95", "АНО «ФИЗКУЛЬТУРА»", "Парковые забеги"],
            ["Ярославское беговое", "ООО «ЯРО. БЕГОВОЕ СООБЩЕСТВО»", "Региональные забеги"],
        ],
        col_widths=[Inches(2.5), Inches(3.5), Inches(3.0)],
    )

    # 4. Динамика рынка
    add_slide(prs, "Динамика рынка: рост в 3.2 раза за 4 года", [
        "Суммарная выручка: 936 → 2 988 млн руб.",
        "Медианный темп роста: +34% г/г",
        "Темпы замедляются: 50% → 20%",
        "Все 7 компаний с данными показали рост",
    ], PLOTS / "market_summary.png", Inches(5.2))

    # 5. Динамика по компаниям
    add_slide(prs, "Динамика ключевых показателей по компаниям",
              image_path=PLOTS / "dynamics.png", image_width=Inches(9.0))

    # 6. Boxplot
    add_slide(prs, "Распределение показателей по компаниям",
              image_path=PLOTS / "boxplots.png", image_width=Inches(8.5))

    # 7. Результаты проверки гипотез
    add_table_slide(prs, "Результаты проверки гипотез",
        ["Гипотеза", "Результат", "p-value", "Ключевой вывод"],
        [
            ["H1: Устойчивый рост", "ПОДТВЕРЖДЕНА", "< 0.0001", "Медиана +34% г/г, выручка ×3.2"],
            ["H2: Масштаб → рентаб.", "Не подтверждена", "0.638", "ρ = −0.08, связь незначима"],
            ["H3: Эффект масштаба", "Не подтверждена", "0.254", "ρ = −0.20, незначима"],
            ["H4: Кластеры", "Не подтв. форм.", "силуэт 0.21", "3 содержательных кластера"],
            ["H5: Рост → рост КЗ", "Не подтверждена", "0.114", "ρ = +0.32, тенденция"],
        ],
        col_widths=[Inches(2.2), Inches(1.8), Inches(1.2), Inches(3.8)],
    )

    # 8. Корреляционная матрица
    add_slide(prs, "Корреляционная матрица (Спирмен)", [
        "Выручка не коррелирует с прибылью (ρ = 0.11)",
        "Активы ↔ КЗ: ρ = +0.92 (рост за счёт КЗ)",
        "Валовая рентабельность → чистая прибыль: ρ = +0.75",
        "Оборачиваемость ↔ активы: ρ = −0.75",
    ], PLOTS / "correlation_heatmap.png", Inches(4.8))

    # 9. Scatter гипотез
    add_slide(prs, "Визуализация гипотез H2, H3, H5",
              image_path=PLOTS / "hypothesis_scatters.png", image_width=Inches(9.0))

    # 10. Регрессия
    add_slide(prs, "Регрессионный анализ", [
        "Модель: Рентабельность = −4.2 + 1.4 · ln(Выручка)",
        "R² = 0.031, p = 0.301 — модель незначима",
        "Условия МНК выполнены (остатки нормальны)",
        "",
        "Рентабельность определяется не масштабом,",
        "а бизнес-моделью и управлением затратами",
    ], PLOTS / "regression.png", Inches(5.0))

    # 11. Кластеры
    add_slide(prs, "Кластерный анализ: 3 типа компаний", [
        "!Кластер 0 — Малые организаторы",
        "Паркран, Петерб. спорт, Арена Плюс, TIMERMAN",
        "Низкая рентабельность (~5%), низкая КЗ",
        "",
        "!Кластер 1 — Кэш-генераторы",
        "IRONSTAR, ГРУТ",
        "Высокая ликвидность (~78% ДС), модель предоплат",
        "",
        "!Кластер 2 — Лидеры рынка",
        "Лига героев, Марафон Сервис",
        "Наибольшая выручка, рентабельность ~16%",
    ], PLOTS / "clusters.png", Inches(5.0))

    # 12. Ключевые выводы
    add_slide(prs, "Ключевые выводы", [
        "+Рынок быстро растёт: суммарная выручка ×3.2 за 4 года, медиана +34% г/г",
        "",
        "-Прибыльность НЕ связана с масштабом — крупнейшая компания убыточна 4 года из 5",
        "",
        "Две модели финансирования: предоплаты (высокая ликвидность) vs кредиторская задолженность",
        "",
        "Себестоимость — ключевой фактор: медианная доля 89%, маржа ~10%",
        "",
        "Бизнес low-asset: компании с минимальными ОС растут быстрее всех",
    ])

    # 13. Рекомендации
    add_slide(prs, "Рекомендации для ЛПР", [
        "!Для организаторов мероприятий",
        "Контролировать себестоимость (маржа ~10%)",
        "Использовать модель предоплат от участников",
        "Не гнаться за масштабом ради масштаба",
        "",
        "!Для инвесторов",
        "Рынок растущий (+34% г/г), входной барьер низкий (low-asset)",
        "Дифференциация бизнес-модели важнее размера",
        "",
        "!Для регуляторов",
        "Рынок прозрачен через ГИРБО",
        "Стандартизировать отчётность НКО для сопоставимости",
    ])

    # Save
    out_path = OUT / "presentation.pptx"
    prs.save(out_path)
    print(f"Готово: {out_path}")
    print(f"Слайдов: {len(prs.slides)}")
    print(f"Размер: {out_path.stat().st_size / 1024:.0f} КБ")


if __name__ == "__main__":
    main()
